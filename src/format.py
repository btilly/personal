#! /usr/bin/env python3

import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

def parse_markup(text):
    """Parses simple markup into a list of slides."""
    slides = []
    current_slide = None

    for line in text.strip().splitlines():
        line = line.rstrip()
        if line.startswith('# '):
            if current_slide:
                slides.append(current_slide)
            current_slide = {'name': f"slide_{len(slides)+1:02d}",
                    'title': line[2:].strip(), 'bullets': []}
        elif line.strip().startswith('-'):
            indent = len(line) - len(line.lstrip())
            bullet = line.strip()[1:].strip()
            current_slide['bullets'].append((indent, bullet))
    if current_slide:
        slides.append(current_slide)
    return slides


def format_run(paragraph, text):
    """Formats a run with *italic* or **bold** markers."""
    run = paragraph.add_run()
    run.font.size = Pt(28)
    # Detect bold + italic
    if re.match(r"\*\*\*(.+?)\*\*\*", text):
        run.text = re.findall(r"\*\*\*(.+?)\*\*\*", text)[0]
        run.font.bold = True
        run.font.italic = True
    elif re.match(r"\*\*(.+?)\*\*", text):
        run.text = re.findall(r"\*\*(.+?)\*\*", text)[0]
        run.font.bold = True
    elif re.match(r"\*(.+?)\*", text):
        run.text = re.findall(r"\*(.+?)\*", text)[0]
        run.font.italic = True
    else:
        run.text = text


def create_pptx_from_markup(text, output_file="output.pptx"):
    slides_data = parse_markup(text)
    prs = Presentation()

    print(prs.slide_layouts[1].shapes.__repr__())

    for slide in slides_data:
        slide_layout = prs.slide_layouts[1]
        ppt_slide = prs.slides.add_slide(slide_layout)
        ppt_slide.name = slide['name']

        title = ppt_slide.shapes.title
        title.text = slide['title'].replace("\\n", "\n")

        content = ppt_slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

        p = text_frame.paragraphs[0]
        i = 0

        for indent, bullet in slide['bullets']:
            if 0 < i:
                p = text_frame.add_paragraph()
            i += 1
            p.level = indent // 2  # every 2 spaces = 1 indent level
            format_run(p, bullet)

        
        #ppt_slide.shapes[0].top = ppt_slide.shapes[0].top - Pt(1)
        #print(ppt_slide.shapes[0].top)
        #ppt_slide.shapes[1].top = ppt_slide.shapes[1].top - Pt(1)

    prs.save(output_file)
    print(f"Saved to {output_file}")

try:
    with open("input.txt", 'r') as file:
        file_content = file.read()
    create_pptx_from_markup(file_content)
except FileNotFoundError:
    print(f"Error: The file '{file_path}' was not found.")
except Exception as e:
    print(f"An error occurred: {e}")

